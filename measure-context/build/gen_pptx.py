# -*- coding: utf-8 -*-
"""PowerPoint (.pptx) renderer with python-pptx.

Unlike openpyxl/python-docx, python-pptx exposes a real drawing API, so the four
diagrams are drawn with native autoshapes (rectangles, rounded rects, diamonds,
folded-corner notes) and straight connectors with arrowheads -- laid out from the
shared layout.py geometry, scaled to the slide. The two captures are embedded as
real pictures (add_picture). Tables use native PowerPoint tables.

build() returns (path, logical_text) where logical_text is every string written
to a shape / textbox / table cell -- the (a) logical-content measure.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

import content as C
import layout as LY
from common import dwidth
from gen_text import SAMPLES, _meta_lines

_INK = RGBColor(0x33, 0x41, 0x5C)
_FILL = {
    "rect": RGBColor(0xDA, 0xE8, 0xFC),
    "roundrect": RGBColor(0xD5, 0xE8, 0xD4),
    "stadium": RGBColor(0xF8, 0xCE, 0xCC),
    "diamond": RGBColor(0xFF, 0xE6, 0xCC),
    "note": RGBColor(0xFF, 0xF2, 0xCC),
}
_MSO = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "stadium": MSO_SHAPE.ROUNDED_RECTANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "note": MSO_SHAPE.FOLDED_CORNER,
}
AREA_L, AREA_T = Inches(0.3), Inches(1.15)
AREA_W, AREA_H = Inches(12.7), Inches(6.0)


def _arrow(conn):
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def _dash(conn):
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.L = []

    def slide(self):
        return self.prs.slides.add_slide(self.blank)

    def titlebox(self, slide, text, size=24):
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = text
        tf.paragraphs[0].font.size = Pt(size)
        tf.paragraphs[0].font.bold = True
        self.L.append(text)

    def textbox(self, slide, lines, left, top, width, height, size=14):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln
            p.font.size = Pt(size)
            p.space_after = Pt(6)
            self.L.append(ln)

    def diagram(self, slide, g, fontpt):
        W, H = g["size"]
        scale = min(AREA_W / W, AREA_H / H)
        X = lambda px: Emu(int(AREA_L + px * scale))
        Y = lambda px: Emu(int(AREA_T + px * scale))
        S = lambda px: Emu(int(px * scale))
        for e in g["edges"]:
            pts = e["points"]
            for i in range(len(pts) - 1):
                a, b = pts[i], pts[i + 1]
                conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                                  X(a[0]), Y(a[1]), X(b[0]), Y(b[1]))
                conn.line.color.rgb = _INK
                conn.line.width = Pt(1)
                if e["dashed"]:
                    _dash(conn)
                if e["arrow"] and i == len(pts) - 2:
                    _arrow(conn)
            if e["label"]:
                i = len(pts) // 2
                mx = (pts[i - 1][0] + pts[i][0]) / 2
                my = (pts[i - 1][1] + pts[i][1]) / 2
                tb = slide.shapes.add_textbox(Emu(int(AREA_L + mx * scale) - Inches(0.6)),
                                              Emu(int(AREA_T + my * scale) - Inches(0.1)),
                                              Inches(1.2), Inches(0.2))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.text = e["label"]
                tf.paragraphs[0].font.size = Pt(8)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                self.L.append(e["label"])
        for s in g["shapes"]:
            shp = slide.shapes.add_shape(_MSO[s["kind"]], X(s["x"]), Y(s["y"]), S(s["w"]), S(s["h"]))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _FILL[s["kind"]]
            shp.line.color.rgb = _INK
            tf = shp.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.text = s["label"]
            p = tf.paragraphs[0]
            p.font.size = Pt(fontpt)
            p.font.color.rgb = _INK
            p.alignment = PP_ALIGN.CENTER
            self.L.append(s["label"])

    def table(self, slide, header, rows, fontpt=10):
        nr, nc = len(rows) + 1, len(header)
        gt = slide.shapes.add_table(nr, nc, Inches(0.3), Inches(1.2),
                                    Inches(12.7), Inches(0.3 * nr)).table
        for c, h in enumerate(header):
            gt.cell(0, c).text = str(h)
            self.L.append(str(h))
        for ri, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                gt.cell(ri, c).text = str(val)
                self.L.append(str(val))
        for r in range(nr):
            for c in range(nc):
                for p in gt.cell(r, c).text_frame.paragraphs:
                    p.font.size = Pt(fontpt)

    def picture(self, slide, path):
        slide.shapes.add_picture(path, Inches(1.8), Inches(1.3), width=Inches(9.7))


def build():
    d = Deck()

    # title slide
    s = d.slide()
    d.titlebox(s, C.DOC["title"], size=28)
    d.textbox(s, [" / ".join("%s: %s" % kv for kv in _meta_lines()),
                  "", "概要: " + C.DOC["abstract"]],
              Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.0), size=16)

    def section_slide(num):
        n, head, paras = next(x for x in C.SECTIONS if x[0] == num)
        sl = d.slide()
        d.titlebox(sl, "%s %s" % (n, head))
        d.textbox(sl, list(paras), Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0), size=13)

    def diagram_slide(caption, key, fontpt):
        sl = d.slide()
        d.titlebox(sl, caption)
        d.diagram(sl, LY.ALL[key](), fontpt)

    section_slide("1")
    diagram_slide(C.FIG_BLOCK, "block", 11)
    section_slide("2")
    section_slide("3")
    diagram_slide(C.FIG_STATE, "state", 11)
    section_slide("4")
    diagram_slide(C.FIG_SEQ, "sequence", 9)
    section_slide("5")
    diagram_slide(C.FIG_FLOW, "flow", 8)
    section_slide("6")
    section_slide("7")

    s = d.slide()
    d.titlebox(s, C.PARAM_TABLE_CAPTION)
    d.table(s, C.PARAM_TABLE_HEADER, C.PARAM_TABLE_ROWS, fontpt=10)
    s = d.slide()
    d.titlebox(s, C.IF_TABLE_CAPTION)
    d.table(s, C.IF_TABLE_HEADER, C.IF_TABLE_ROWS, fontpt=9)

    for cap, img in [(C.FIG_CAP1, C.CAPTURE1), (C.FIG_CAP2, C.CAPTURE2)]:
        s = d.slide()
        d.titlebox(s, cap)
        d.picture(s, os.path.join(SAMPLES, img))

    path = os.path.join(SAMPLES, "spec_lkas.pptx")
    d.prs.save(path)
    return path, "\n".join(d.L)


if __name__ == "__main__":
    from common import count_tokens
    p, logical = build()
    print("pptx logical:", count_tokens(logical), "tokens,", len(logical), "chars ->", os.path.basename(p))
